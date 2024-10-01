#!/usr/bin/env python
# -*- coding: utf-8 -*-

import os
from typing import Dict, List, Optional, Tuple, Union
import PyPDF2
import markdown
import html2text
import json
from tqdm import tqdm
import tiktoken
from bs4 import BeautifulSoup
import re
import pptx
from langchain_community.document_loaders import PyMuPDFLoader
import docx
from pptx import Presentation

enc = tiktoken.get_encoding("cl100k_base")


class ReadFiles:
    """
    class to read files
    """

    def __init__(self, path: str) -> None:
        self._path = path
        self.file_list = self.get_files()
        self.file_list_single_folder = self.get_files_single_folder()

    def get_files_single_folder(self):
        file_list = []
        # 列出self._path目录下的所有内容
        for entry in os.listdir(self._path):
            # 获取文件的完整路径
            full_path = os.path.join(self._path, entry)
            # 检查是否是文件，并且后缀名符合要求
            if os.path.isfile(full_path) and os.path.splitext(entry)[1] in {".md", ".txt", ".pdf", ".docx", ".pptx", ".py", ".cpp", ".json"}:
                file_list.append(full_path)
        return file_list

    def get_files(self):
        # args：dir_path，目标文件夹路径
        file_list = []
        for filepath, dirnames, filenames in os.walk(self._path):
            # os.walk 函数将递归遍历指定文件夹
            for filename in filenames:
                # 通过后缀名判断文件类型是否满足要求
                if filename.endswith(".md"):
                    # 如果满足要求，将其绝对路径加入到结果列表
                    file_list.append(os.path.join(filepath, filename))
                elif filename.endswith(".txt"):
                    file_list.append(os.path.join(filepath, filename))
                elif filename.endswith(".pdf"):
                    file_list.append(os.path.join(filepath, filename))
                elif filename.endswith(".docx"):
                    file_list.append(os.path.join(filepath, filename))
                elif filename.endswith(".pptx"):
                    file_list.append(os.path.join(filepath, filename))
                elif filename.endswith(".py"):
                    file_list.append(os.path.join(filepath, filename))
                elif filename.endswith(".cpp"):
                    file_list.append(os.path.join(filepath, filename))
                elif filename.endswith(".json"):
                    file_list.append(os.path.join(filepath, filename))
        return file_list

    def get_content_single_folder(self, max_token_len: int = 1000, cover_content: int = 150):
        docs = []
        # 读取文件内容
        for file in self.file_list_single_folder:
            content = self.read_file_content(file)
            content = self.process_text(content)
            chunk_content = self.get_chunk(
                content, max_token_len=max_token_len, cover_content=cover_content)
            docs.extend(chunk_content)
        return docs

    def process_text(self,text):
        start_index = text.find("【爸比偏好观察】")
        end_index = text.find("【爸比心情观察】")
        if start_index != -1 and end_index != -1:
            text = text[:start_index] + text[end_index + len("【爸比心情观察】"):]

        start_index = text.find("【社交关系】")
        end_index = text.find("【漆小喵的感想】")
        if start_index != -1 and end_index != -1:
            text = text[:start_index] + text[end_index + len("【漆小喵的感想】"):]
        
        text = text.replace("【日期】","")
        text = text.replace("【摘要】","")
        text = text.replace("【爸比行为观察】","")
        text = text.replace("【爸比这段时间要做的事】","")

        return text

    def get_content(self, max_token_len: int = 1000, cover_content: int = 150):
        docs = []
        # 读取文件内容
        for file in self.file_list:
            content = self.read_file_content(file)
            content = self.process_text(content)
            chunk_content = self.get_chunk(
                content, max_token_len=max_token_len, cover_content=cover_content)
            docs.extend(chunk_content)
        return docs

    @classmethod
    def get_chunk(cls, text: str, max_token_len: int = 600, cover_content: int = 150):
        chunk_text = []

        curr_len = 0
        curr_chunk = ''

        token_len = max_token_len - cover_content
        lines = text.splitlines()  # 假设以换行符分割文本为行

        for line in lines:
            line = line.replace(' ', '')
            line_len = len(enc.encode(line))
            if line_len > max_token_len:
                # 如果单行长度就超过限制，则将其分割成多个块
                num_chunks = (line_len + token_len - 1) // token_len
                for i in range(num_chunks):
                    start = i * token_len
                    end = start + token_len
                    # 避免跨单词分割
                    while not line[start:end].rstrip().isspace():
                        start += 1
                        end += 1
                        if start >= line_len:
                            break
                    curr_chunk = curr_chunk[-cover_content:] + line[start:end]
                    chunk_text.append(curr_chunk)
                # 处理最后一个块
                start = (num_chunks - 1) * token_len
                curr_chunk = curr_chunk[-cover_content:] + line[start:end]
                chunk_text.append(curr_chunk)
                
            if curr_len + line_len <= token_len:
                curr_chunk += line
                curr_chunk += '\n'
                curr_len += line_len
                curr_len += 1
            else:
                chunk_text.append(curr_chunk)
                curr_chunk = curr_chunk[-cover_content:]+line
                curr_len = line_len + cover_content

        if curr_chunk:
            chunk_text.append(curr_chunk)

        return chunk_text

    # @classmethod
    # def read_file_content(cls, file_path: str):
    #     # 根据文件扩展名选择读取方法
    #     if file_path.endswith('.pdf'):
    #         return cls.read_pdf(file_path)
    #     elif file_path.endswith('.md'):
    #         return cls.read_markdown(file_path)
    #     elif file_path.endswith('.txt'):
    #         return cls.read_text(file_path)
    #     elif 
    #     else:
    #         raise ValueError("Unsupported file type")

# class FileReader:
    @classmethod
    def read_file_content(cls, file_path: str):
        # 根据文件扩展名选择读取方法
        if file_path.endswith('.pdf'):
            return cls.read_pdf(file_path)
        elif file_path.endswith('.md'):
            return cls.read_markdown(file_path)
        elif file_path.endswith('.txt'):
            return cls.read_text(file_path)
        elif file_path.endswith('.py'):
            return cls.read_text(file_path)
        elif file_path.endswith('.docx'):
            return cls.read_docx(file_path)
        elif file_path.endswith('.pptx'):
            return cls.read_pptx(file_path)
        elif file_path.endswith('.json'):
            return cls.read_json(file_path)
        elif file_path.endswith('.cpp'):
            return cls.read_text(file_path)
        else:
            raise ValueError("Unsupported file type")

    @classmethod
    def read_pdf(cls, file_path: str):
        # 读取PDF文件
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page_num in range(len(reader.pages)):
                text += reader.pages[page_num].extract_text()
            return text

    @classmethod
    def read_markdown(cls, file_path: str):
        # 读取Markdown文件
        with open(file_path, 'r', encoding='utf-8') as file:
            md_text = file.read()
            html_text = markdown.markdown(md_text)
            # 使用BeautifulSoup从HTML中提取纯文本
            soup = BeautifulSoup(html_text, 'html.parser')
            plain_text = soup.get_text()
            # 使用正则表达式移除网址链接
            text = re.sub(r'http\S+', '', plain_text) 
            return text

    @classmethod
    def read_text(cls, file_path: str):
        # 读取文本文件
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    @classmethod
    def read_docx(cls, file_path: str):
        # 读取docx文件
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)

    @classmethod
    def read_pptx(cls, file_path: str):
        # 读取pptx文件
        presentation = pptx.Presentation(file_path)
        full_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)

    @classmethod
    def read_json(cls, file_path: str):
        # 读取json文件
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            return json.dumps(data, indent=4, ensure_ascii=False)


class Documents:
    """
        获取已分好类的json格式文档
    """
    def __init__(self, path: str = '') -> None:
        self.path = path
    
    def get_content(self):
        with open(self.path, mode='r', encoding='utf-8') as f:
            content = json.load(f)
        return content


def extract_text(path):
    return open(path, 'r',encoding = "utf-8").read()

def extract_pdf(path):
    loader = PyMuPDFLoader(path)
    data = loader.load()
    data = [x.page_content for x in data]
    content = '\n\n'.join(data)
    return content

def extract_docx(path):
    doc = docx.Document(path)
    data = []
    for paragraph in doc.paragraphs:
        data.append(paragraph.text)
    content = '\n\n'.join(data)
    return content

def extract_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text